from pptx import Presentation


def extract_notes(pptx_file):
    # 加载PPT文件
    presentation = Presentation(pptx_file)

    # 用来存储备注信息的字符串
    notes_text = ""

    # 遍历每一页幻灯片
    for slide_num, slide in enumerate(presentation.slides):
        # 检查是否有备注
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes = notes_slide.notes_text_frame.text
            # 将备注按页格式化并添加到文本中
            notes_text += f"page{slide_num + 1}: {notes}\n"
        else:
            notes_text += f"page{slide_num + 1}: No notes\n"

    return notes_text


# 使用函数来提取PPT的备注信息
pptx_file = r"F:\项目\2025竞赛\科技活动专项项目结题验收-武汉大学-潘俊0315修订珊瑚、单木和竞赛平台.pptx"  # 替换为你的PPT文件路径
notes = extract_notes(pptx_file)

# 输出备注信息
print(notes)
